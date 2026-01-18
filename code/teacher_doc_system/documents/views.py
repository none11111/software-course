from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.mixins import LoginRequiredMixin, UserPassesTestMixin
from django.contrib import messages
from django.urls import reverse_lazy, reverse
from django.views.generic import View, ListView, CreateView, UpdateView, DeleteView, DetailView, TemplateView
from django.http import JsonResponse, HttpResponse, Http404, FileResponse
from django.db.models import Q, Count, Sum
from django.utils import timezone
from django.core.paginator import Paginator
from django.db import transaction
from django.conf import settings
import os
import hashlib
import mimetypes
from datetime import datetime, timedelta

from .models import Document, DocumentCategory, DocumentVersion, DocumentOperationLog
from system.models import ShareLink
from .forms import DocumentForm, CategoryForm, VersionForm, ShareLinkForm
from django.contrib.auth import get_user_model

User = get_user_model()


class TeacherDashboardView(LoginRequiredMixin, TemplateView):
    """教师仪表盘"""
    template_name = 'documents/teacher_dashboard.html'
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        user = self.request.user
        
        # 基本统计
        total_documents = Document.objects.filter(author=user).count()
        public_documents = Document.objects.filter(author=user, is_public=True).count()
        private_documents = total_documents - public_documents
        
        # 按状态统计文档
        status_stats = Document.objects.filter(author=user).values('status').annotate(
            count=Count('id')
        )
        
        # 审核未通过的文档
        rejected_documents = Document.objects.filter(author=user, status='rejected').count()
        
        # 存储使用情况
        total_storage_used = user.storage_used
        total_storage_quota = user.storage_quota
        
        # 计算存储使用百分比，确保显示至少2位小数
        storage_percentage = 0
        if total_storage_quota > 0:
            percentage = (total_storage_used / total_storage_quota) * 100
            # 如果百分比很小，至少显示0.01%
            if percentage > 0 and percentage < 0.01:
                storage_percentage = 0.01
            else:
                storage_percentage = round(percentage, 2)
        
        # 按分类统计文档
        category_stats = Document.objects.filter(author=user).values(
            'category__name'
        ).annotate(
            count=Count('id')
        ).order_by('-count')[:5]
        
        # 按文件类型统计
        file_type_stats = Document.objects.filter(author=user).values(
            'file_type'
        ).annotate(
            count=Count('id')
        ).order_by('-count')[:5]
        
        # 最近上传的文档
        recent_documents = Document.objects.filter(author=user).order_by('-created_at')[:5]
        
        # 分享链接统计
        total_share_links = ShareLink.objects.filter(created_by=user).count()
        active_share_links = ShareLink.objects.filter(
            created_by=user, 
            is_active=True,
            expires_at__gt=timezone.now()
        ).count()
        
        # 最近7天的活动统计
        week_ago = timezone.now() - timedelta(days=7)
        recent_uploads = Document.objects.filter(
            author=user, 
            created_at__gte=week_ago
        ).count()
        
        recent_operations = DocumentOperationLog.objects.filter(
            user=user,
            created_at__gte=week_ago
        ).count()
        
        # 存储使用趋势（最近7天）
        storage_trend = []
        for i in range(7):
            date = timezone.now() - timedelta(days=i)
            day_start = date.replace(hour=0, minute=0, second=0, microsecond=0)
            day_end = day_start + timedelta(days=1)
            
            # 计算到该日期为止的存储使用量
            documents_until_date = Document.objects.filter(
                author=user,
                created_at__lt=day_end
            ).aggregate(total_size=Sum('file_size'))['total_size'] or 0
            
            storage_trend.append({
                'date': day_start.strftime('%m-%d'),
                'size': documents_until_date
            })
        
        storage_trend.reverse()  # 按时间正序排列
        
        context.update({
            'total_documents': total_documents,
            'public_documents': public_documents,
            'private_documents': private_documents,
            'rejected_documents': rejected_documents,
            'status_stats': status_stats,
            'total_storage_used': total_storage_used,
            'total_storage_quota': total_storage_quota,
            'storage_percentage': storage_percentage,
            'category_stats': category_stats,
            'file_type_stats': file_type_stats,
            'recent_documents': recent_documents,
            'total_share_links': total_share_links,
            'active_share_links': active_share_links,
            'recent_uploads': recent_uploads,
            'recent_operations': recent_operations,
            'storage_trend': storage_trend,
        })
        
        return context


class DocumentListView(LoginRequiredMixin, ListView):
    """文档列表"""
    model = Document
    template_name = 'documents/document_list.html'
    context_object_name = 'documents'
    paginate_by = 20
    
    def get_queryset(self):
        user = self.request.user
        
        # 基础查询：自己的文档 + 公开文档
        queryset = Document.objects.filter(
            Q(author=user) | Q(is_public=True)
        ).select_related('author', 'category')
        
        # 筛选条件
        category_filter = self.request.GET.get('category')
        status_filter = self.request.GET.get('status')
        file_type_filter = self.request.GET.get('file_type')
        date_from = self.request.GET.get('date_from')
        date_to = self.request.GET.get('date_to')
        search = self.request.GET.get('search')
        my_docs_only = self.request.GET.get('my_docs_only')
        
        if my_docs_only == 'true':
            queryset = queryset.filter(author=user)
        
        if category_filter:
            queryset = queryset.filter(category=category_filter)
        
        if status_filter:
            queryset = queryset.filter(status=status_filter)
        
        if file_type_filter:
            queryset = queryset.filter(file_type=file_type_filter)
        
        if date_from:
            from datetime import datetime
            try:
                # 确保日期格式正确
                if isinstance(date_from, str):
                    date_from = datetime.strptime(date_from, '%Y-%m-%d').date()
                # 使用更兼容的日期筛选方法
                from django.utils import timezone
                start_datetime = timezone.datetime.combine(date_from, timezone.datetime.min.time())
                start_datetime = timezone.make_aware(start_datetime)
                queryset = queryset.filter(created_at__gte=start_datetime)
            except ValueError:
                # 如果日期格式不正确，忽略该筛选条件
                pass
        
        if date_to:
            from datetime import datetime
            try:
                # 确保日期格式正确
                if isinstance(date_to, str):
                    date_to = datetime.strptime(date_to, '%Y-%m-%d').date()
                # 使用更兼容的日期筛选方法
                from django.utils import timezone
                end_datetime = timezone.datetime.combine(date_to, timezone.datetime.max.time())
                end_datetime = timezone.make_aware(end_datetime)
                queryset = queryset.filter(created_at__lte=end_datetime)
            except ValueError:
                # 如果日期格式不正确，忽略该筛选条件
                pass
        
        if search:
            queryset = queryset.filter(
                Q(title__icontains=search) |
                Q(description__icontains=search) |
                Q(author__first_name__icontains=search) |
                Q(author__last_name__icontains=search)
            )
        
        return queryset.order_by('-created_at')
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        
        # 添加筛选选项 - 根据用户权限限制分类选择
        user = self.request.user
        if user.is_superuser or user.is_admin():
            # 管理员可以看到所有分类
            categories = DocumentCategory.objects.filter(is_active=True)
        else:
            # 教师只能看到管理员创建的分类和自己创建的分类
            admin_users = User.objects.filter(Q(is_superuser=True) | Q(role='admin'))
            categories = DocumentCategory.objects.filter(
                Q(created_by__in=admin_users) | Q(created_by=user),
                is_active=True
            )
        
        file_types = Document.objects.values_list('file_type', flat=True).distinct()
        
        context.update({
            'categories': categories,
            'file_types': file_types,
            'status_choices': Document.STATUS_CHOICES,
            'filters': {
                'category': self.request.GET.get('category'),
                'status': self.request.GET.get('status'),
                'file_type': self.request.GET.get('file_type'),
                'date_from': self.request.GET.get('date_from'),
                'date_to': self.request.GET.get('date_to'),
                'search': self.request.GET.get('search'),
                'my_docs_only': self.request.GET.get('my_docs_only'),
            }
        })
        
        return context


class DocumentSearchView(LoginRequiredMixin, ListView):
    """文档搜索"""
    model = Document
    template_name = 'documents/document_list.html'
    context_object_name = 'documents'
    paginate_by = 20
    
    def get_queryset(self):
        user = self.request.user
        
        # 基础查询：自己的文档 + 公开文档
        queryset = Document.objects.filter(
            Q(author=user) | Q(is_public=True)
        ).select_related('author', 'category')
        
        # 搜索关键词
        search = self.request.GET.get('search')
        if search:
            queryset = queryset.filter(
                Q(title__icontains=search) |
                Q(description__icontains=search) |
                Q(author__username__icontains=search)
            )
        
        # 其他筛选条件
        category_filter = self.request.GET.get('category')
        status_filter = self.request.GET.get('status')
        file_type_filter = self.request.GET.get('file_type')
        date_from = self.request.GET.get('date_from')
        date_to = self.request.GET.get('date_to')
        my_docs_only = self.request.GET.get('my_docs_only')
        
        if my_docs_only == 'true':
            queryset = queryset.filter(author=user)
        
        if category_filter:
            queryset = queryset.filter(category=category_filter)
        
        if status_filter:
            queryset = queryset.filter(status=status_filter)
        
        if file_type_filter:
            queryset = queryset.filter(file_type=file_type_filter)
        
        if date_from:
            from datetime import datetime
            try:
                # 确保日期格式正确
                if isinstance(date_from, str):
                    date_from = datetime.strptime(date_from, '%Y-%m-%d').date()
                # 使用更兼容的日期筛选方法
                from django.utils import timezone
                start_datetime = timezone.datetime.combine(date_from, timezone.datetime.min.time())
                start_datetime = timezone.make_aware(start_datetime)
                queryset = queryset.filter(created_at__gte=start_datetime)
            except ValueError:
                # 如果日期格式不正确，忽略该筛选条件
                pass
        
        if date_to:
            from datetime import datetime
            try:
                # 确保日期格式正确
                if isinstance(date_to, str):
                    date_to = datetime.strptime(date_to, '%Y-%m-%d').date()
                # 使用更兼容的日期筛选方法
                from django.utils import timezone
                end_datetime = timezone.datetime.combine(date_to, timezone.datetime.max.time())
                end_datetime = timezone.make_aware(end_datetime)
                queryset = queryset.filter(created_at__lte=end_datetime)
            except ValueError:
                # 如果日期格式不正确，忽略该筛选条件
                pass
        
        return queryset.order_by('-created_at')
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        
        # 根据用户权限限制分类选择
        user = self.request.user
        if user.is_superuser or user.is_admin():
            # 管理员可以看到所有分类
            categories = DocumentCategory.objects.filter(is_active=True)
        else:
            # 教师只能看到管理员创建的分类和自己创建的分类
            admin_users = User.objects.filter(Q(is_superuser=True) | Q(role='admin'))
            categories = DocumentCategory.objects.filter(
                Q(created_by__in=admin_users) | Q(created_by=user),
                is_active=True
            )
        
        context['categories'] = categories
        context['search_form'] = self.get_search_form()
        return context
    
    def get_search_form(self):
        from .forms import DocumentSearchForm
        return DocumentSearchForm(self.request.GET)


class DocumentDetailView(LoginRequiredMixin, DetailView):
    """文档详情"""
    model = Document
    template_name = 'documents/document_detail.html'
    context_object_name = 'document'
    
    def get_queryset(self):
        user = self.request.user
        # 管理员不能访问文档详情页面
        if user.is_superuser or user.is_admin():
            return Document.objects.none()
        # 普通用户只能访问自己的文档或公开文档
        return Document.objects.filter(
            Q(author=user) | Q(is_public=True)
        ).select_related('author', 'category')
    
    def get_context_data(self, **kwargs):
        """添加上下文数据，处理文件路径问题"""
        context = super().get_context_data(**kwargs)
        document = context['document']
        
        # 处理文件路径问题（旧路径可能包含重复的 media/）
        if document.file:
            try:
                # 尝试访问文件大小，如果失败则使用存储的 file_size
                _ = document.file.size
            except (ValueError, OSError, FileNotFoundError):
                # 如果文件路径有问题，使用存储的 file_size 而不是访问文件
                # 这样模板中访问 document.file.size 时不会报错
                # 但我们需要在模板中使用 document.file_size 而不是 document.file.size
                pass
        
        return context
    
    def get(self, request, *args, **kwargs):
        # 如果是管理员，重定向到审核页面
        if request.user.is_superuser or request.user.is_admin():
            document = get_object_or_404(Document, pk=kwargs['pk'])
            if document.status == 'review':
                return redirect('documents:document_review', pk=document.pk)
            else:
                messages.info(request, '管理员请通过文档审核页面管理文档')
                return redirect('documents:document_review_list')
        
        document = self.get_object()
        
        # 增加查看次数
        if document.author != request.user:
            document.view_count += 1
            document.save()
        
        # 记录查看日志
        DocumentOperationLog.objects.create(
            document=document,
            user=request.user,
            operation='view',
            ip_address=self._get_client_ip(),
            details={}
        )
        
        return super().get(request, *args, **kwargs)
    
    def _get_client_ip(self):
        """获取客户端IP地址"""
        x_forwarded_for = self.request.META.get('HTTP_X_FORWARDED_FOR')
        if x_forwarded_for:
            ip = x_forwarded_for.split(',')[0]
        else:
            ip = self.request.META.get('REMOTE_ADDR')
        return ip


class UploadDocumentView(LoginRequiredMixin, CreateView):
    """上传文档"""
    model = Document
    form_class = DocumentForm
    template_name = 'documents/upload_document.html'
    success_url = reverse_lazy('documents:document_list')
    
    def get_form_kwargs(self):
        kwargs = super().get_form_kwargs()
        kwargs['user'] = self.request.user
        return kwargs
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        
        # 根据用户权限获取分类，按层级结构组织
        user = self.request.user
        if user.is_superuser or user.is_admin():
            # 管理员可以看到所有分类
            categories = DocumentCategory.objects.filter(is_active=True).select_related('parent', 'created_by').prefetch_related('children')
        else:
            # 教师只能看到管理员创建的分类和自己创建的分类
            admin_users = User.objects.filter(Q(is_superuser=True) | Q(role='admin'))
            categories = DocumentCategory.objects.filter(
                Q(created_by__in=admin_users) | Q(created_by=user),
                is_active=True
            ).select_related('parent', 'created_by').prefetch_related('children')
        
        context['categories'] = categories
        return context
    
    def form_valid(self, form):
        user = self.request.user
        file = form.cleaned_data['file']
        
        # 检查存储配额
        if not user.can_upload_file(file.size):
            messages.error(self.request, '存储空间不足，无法上传文件')
            return self.form_invalid(form)
        
        # 计算文件哈希
        file_hash = self._calculate_file_hash(file)
        
        # 检查文件是否已存在（检查是否已有相同哈希的文档）
        existing_doc = Document.objects.filter(file_hash=file_hash).first()
        if existing_doc:
            # 如果文件已存在，提示用户并允许选择
            # 但为了避免唯一性约束错误，我们需要处理
            # 方案1：允许重复哈希（需要修改模型）
            # 方案2：在哈希后添加时间戳使其唯一
            # 这里采用方案2，在哈希后添加时间戳，保持去重功能的同时避免唯一性错误
            import time
            file_hash = f"{file_hash}_{int(time.time() * 1000000)}"  # 添加微秒时间戳使其唯一
            messages.info(self.request, '检测到相同内容的文件，已创建新的文档记录')
        
        with transaction.atomic():
            # 创建文档
            document = form.save(commit=False)
            document.author = user
            document.file_size = file.size
            document.file_type = self._get_file_type(file)
            document.file_hash = file_hash
            document.save()
            
            # 创建初始版本记录（v1.0）
            DocumentVersion.objects.create(
                document=document,
                version_number='v1.0',
                file=document.file,
                file_size=document.file_size,
                change_log='初始版本',
                created_by=user
            )
            
            # 更新用户存储使用量
            user.storage_used += file.size
            user.save()
            
            # 记录操作日志
            DocumentOperationLog.objects.create(
                document=document,
                user=user,
                operation='create',
                ip_address=self._get_client_ip(),
                details={'file_size': file.size, 'file_type': document.file_type}
            )
            
            messages.success(self.request, f'文档 "{document.title}" 上传成功')
        
        return super().form_valid(form)
    
    def _calculate_file_hash(self, file):
        """计算文件哈希值"""
        file.seek(0)
        content = file.read()
        file.seek(0)
        return hashlib.md5(content).hexdigest()
    
    def _get_file_type(self, file):
        """获取文件类型"""
        ext = os.path.splitext(file.name)[1].lower()
        return ext[1:] if ext else 'unknown'
    
    def _get_client_ip(self):
        """获取客户端IP地址"""
        x_forwarded_for = self.request.META.get('HTTP_X_FORWARDED_FOR')
        if x_forwarded_for:
            ip = x_forwarded_for.split(',')[0]
        else:
            ip = self.request.META.get('REMOTE_ADDR')
        return ip


class EditDocumentView(LoginRequiredMixin, UpdateView):
    """编辑文档"""
    model = Document
    form_class = DocumentForm
    template_name = 'documents/edit_document.html'
    
    def get_queryset(self):
        # 只能编辑自己的文档，包括审核未通过的文档
        return Document.objects.filter(author=self.request.user)
    
    def get_form_kwargs(self):
        kwargs = super().get_form_kwargs()
        kwargs['user'] = self.request.user
        return kwargs
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        
        # 根据用户权限获取分类，按层级结构组织
        user = self.request.user
        if user.is_superuser or user.is_admin():
            # 管理员可以看到所有分类
            categories = DocumentCategory.objects.filter(is_active=True).select_related('parent', 'created_by').prefetch_related('children')
        else:
            # 教师只能看到管理员创建的分类和自己创建的分类
            admin_users = User.objects.filter(Q(is_superuser=True) | Q(role='admin'))
            categories = DocumentCategory.objects.filter(
                Q(created_by__in=admin_users) | Q(created_by=user),
                is_active=True
            ).select_related('parent', 'created_by').prefetch_related('children')
        
        context['categories'] = categories
        return context
    
    def get_success_url(self):
        return reverse_lazy('documents:document_detail', kwargs={'pk': self.object.pk})
    
    def form_valid(self, form):
        old_document = Document.objects.get(pk=self.object.pk)
        
        with transaction.atomic():
            document = form.save()
            
            # 如果文档之前是审核未通过状态，现在重新提交审核
            if old_document.status == 'rejected':
                document.status = 'review'
                document.save()
                messages.success(self.request, f'文档 "{document.title}" 已重新提交审核')
            else:
                messages.success(self.request, f'文档 "{document.title}" 更新成功')
            
            # 记录操作日志
            DocumentOperationLog.objects.create(
                document=document,
                user=self.request.user,
                operation='update',
                ip_address=self._get_client_ip(),
                details={'changes': form.changed_data}
            )
        
        return super().form_valid(form)
    
    def _get_client_ip(self):
        """获取客户端IP地址"""
        x_forwarded_for = self.request.META.get('HTTP_X_FORWARDED_FOR')
        if x_forwarded_for:
            ip = x_forwarded_for.split(',')[0]
        else:
            ip = self.request.META.get('REMOTE_ADDR')
        return ip


class DeleteDocumentView(LoginRequiredMixin, DeleteView):
    """删除文档"""
    model = Document
    template_name = 'documents/document_confirm_delete.html'
    success_url = reverse_lazy('documents:document_list')
    
    def get_queryset(self):
        # 只能删除自己的文档
        return Document.objects.filter(author=self.request.user)
    
    def delete(self, request, *args, **kwargs):
        document = self.get_object()
        
        with transaction.atomic():
            # 更新用户存储使用量
            request.user.storage_used -= document.file_size
            request.user.save()
            
            # 删除文件
            if document.file and os.path.exists(document.file.path):
                os.remove(document.file.path)
            
            # 记录操作日志
            DocumentOperationLog.objects.create(
                document=document,
                user=request.user,
                operation='delete',
                ip_address=self._get_client_ip(),
                details={'file_size': document.file_size}
            )
            
            messages.success(request, f'文档 "{document.title}" 已删除')
        
        return super().delete(request, *args, **kwargs)
    
    def _get_client_ip(self):
        """获取客户端IP地址"""
        x_forwarded_for = self.request.META.get('HTTP_X_FORWARDED_FOR')
        if x_forwarded_for:
            ip = x_forwarded_for.split(',')[0]
        else:
            ip = self.request.META.get('REMOTE_ADDR')
        return ip


class DownloadDocumentView(LoginRequiredMixin, View):
    """下载文档"""
    
    def get(self, request, pk):
        document = get_object_or_404(Document, pk=pk)
        
        # 权限检查：自己的文档、公开文档或管理员
        if (document.author != request.user and 
            not document.is_public and 
            not (request.user.is_superuser or request.user.is_admin())):
            raise Http404("文档不存在")
        
        if not document.file or not os.path.exists(document.file.path):
            messages.error(request, '文件不存在')
            raise Http404("文件不存在")
        
        # 增加下载次数
        document.download_count += 1
        document.save()
        
        # 记录下载日志
        DocumentOperationLog.objects.create(
            document=document,
            user=request.user,
            operation='download',
            ip_address=self._get_client_ip(),
            details={}
        )
        
        # 返回文件
        response = FileResponse(
            document.file,
            content_type=mimetypes.guess_type(document.file.name)[0] or 'application/octet-stream'
        )
        # 强制下载，不预览
        response['Content-Disposition'] = f'attachment; filename="{document.title}{os.path.splitext(document.file.name)[1]}"'
        
        return response
    
    def _get_client_ip(self):
        """获取客户端IP地址"""
        x_forwarded_for = self.request.META.get('HTTP_X_FORWARDED_FOR')
        if x_forwarded_for:
            ip = x_forwarded_for.split(',')[0]
        else:
            ip = self.request.META.get('REMOTE_ADDR')
        return ip


class BatchDeleteDocumentsView(LoginRequiredMixin, View):
    """批量删除当前用户的文档"""
    def post(self, request):
        # 支持两种字段名：document_ids 或 ids（前端可能命名不同）
        id_list = request.POST.getlist('document_ids') or request.POST.getlist('ids')
        # 去重并过滤空值
        try:
            document_ids = list({int(i) for i in id_list if str(i).strip()})
        except ValueError:
            messages.error(request, '参数格式不正确')
            return redirect('documents:document_list')

        if not document_ids:
            messages.info(request, '请选择要删除的文档')
            return redirect('documents:document_list')

        # 仅允许删除自己的文档
        queryset = Document.objects.filter(author=request.user, id__in=document_ids)

        if not queryset.exists():
            messages.info(request, '未找到可删除的文档')
            return redirect('documents:document_list')

        deleted_count = 0
        total_freed = 0

        with transaction.atomic():
            for document in queryset:
                # 备份文件路径与大小
                file_path = document.file.path if document.file else None
                file_size = document.file_size or 0

                # 更新用户存储
                request.user.storage_used = max(0, (request.user.storage_used or 0) - file_size)
                request.user.save()

                # 记录操作日志
                DocumentOperationLog.objects.create(
                    document=document,
                    user=request.user,
                    operation='delete',
                    ip_address=self._get_client_ip(request),
                    details={'file_size': file_size}
                )

                # 删除主文件（忽略不存在）
                try:
                    if file_path and os.path.exists(file_path):
                        os.remove(file_path)
                except Exception:
                    # 文件删除失败不阻断数据库删除
                    pass

                # 删除文档本身（级联删除版本、分享链接等由模型关系处理）
                document.delete()
                deleted_count += 1
                total_freed += file_size

        # 反馈
        if deleted_count:
            messages.success(request, f'已删除 {deleted_count} 个文档，释放存储 {self._format_size(total_freed)}')
        else:
            messages.info(request, '没有文档被删除')

        return redirect('documents:document_list')

    def _format_size(self, size_bytes):
        units = ['B', 'KB', 'MB', 'GB', 'TB']
        size = float(size_bytes)
        idx = 0
        while size >= 1024 and idx < len(units) - 1:
            size /= 1024.0
            idx += 1
        return f"{size:.1f} {units[idx]}"

    def _get_client_ip(self, request):
        x_forwarded_for = request.META.get('HTTP_X_FORWARDED_FOR')
        if x_forwarded_for:
            ip = x_forwarded_for.split(',')[0]
        else:
            ip = request.META.get('REMOTE_ADDR')
        return ip


class CreateShareLinkView(LoginRequiredMixin, View):
    """创建分享链接"""
    template_name = 'documents/create_share_link.html'
    
    def get(self, request, pk):
        """显示创建分享链接表单"""
        document = get_object_or_404(Document, pk=pk, author=request.user)
        form = ShareLinkForm()
        return render(request, self.template_name, {
            'document': document,
            'form': form
        })
    
    def post(self, request, pk):
        """处理分享链接创建"""
        document = get_object_or_404(Document, pk=pk, author=request.user)
        
        form = ShareLinkForm(request.POST)
        if form.is_valid():
            with transaction.atomic():
                # 生成唯一令牌
                token = self._generate_token()
                
                # 创建分享链接
                share_link = ShareLink.objects.create(
                    document=document,
                    token=token,
                    password=form.cleaned_data.get('password', ''),
                    expires_at=form.cleaned_data['expires_at'],
                    max_downloads=form.cleaned_data['max_downloads'],
                    created_by=request.user
                )
                
                messages.success(request, '分享链接创建成功')
                return redirect('documents:share_link_list')
        
        return render(request, self.template_name, {
            'document': document,
            'form': form
        })
    
    def _generate_token(self):
        """生成唯一令牌"""
        import secrets
        return secrets.token_urlsafe(32)


class ShareLinkView(View):
    """分享链接页面"""
    template_name = 'documents/share_link.html'
    
    def get(self, request, token):
        share_link = get_object_or_404(ShareLink, token=token)
        
        if not share_link.is_available:
            return render(request, 'documents/share_link_expired.html', {
                'share_link': share_link
            })
        
        # 计算剩余下载次数
        remaining_downloads = 0
        if share_link.max_downloads > 0:
            remaining_downloads = max(0, share_link.max_downloads - share_link.download_count)
        
        return render(request, self.template_name, {
            'share_link': share_link,
            'document': share_link.document,
            'remaining_downloads': remaining_downloads
        })
    
    def post(self, request, token):
        share_link = get_object_or_404(ShareLink, token=token)
        password = request.POST.get('password', '')
        
        if not share_link.can_access(password):
            messages.error(request, '密码错误或链接已失效')
            return render(request, self.template_name, {
                'share_link': share_link,
                'document': share_link.document
            })
        
        # 密码验证成功，重定向到下载页面
        download_url = reverse('documents:download_share', kwargs={'token': token})
        if password:
            download_url += f'?password={password}'
        return redirect(download_url)


class DownloadShareView(View):
    """下载分享文档"""
    
    def get(self, request, token):
        share_link = get_object_or_404(ShareLink, token=token)
        
        if not share_link.is_available:
            return render(request, 'documents/share_link_expired.html')
        
        # 检查密码
        password = request.GET.get('password', '')
        if share_link.password and not share_link.can_access(password):
            messages.error(request, '密码错误')
            return redirect('documents:share_link', token=token)
        
        # 检查下载次数限制
        if share_link.max_downloads > 0 and share_link.download_count >= share_link.max_downloads:
            messages.error(request, '下载次数已达上限')
            return redirect('documents:share_link', token=token)
        
        if not share_link.document.file or not os.path.exists(share_link.document.file.path):
            messages.error(request, '文件不存在')
            return redirect('documents:share_link', token=token)
        
        # 增加下载次数
        share_link.download_count += 1
        share_link.save()
        
        # 返回文件
        response = FileResponse(
            share_link.document.file,
            content_type=mimetypes.guess_type(share_link.document.file.name)[0] or 'application/octet-stream'
        )
        # 强制下载，不预览
        response['Content-Disposition'] = f'attachment; filename="{share_link.document.title}{os.path.splitext(share_link.document.file.name)[1]}"'
        
        return response


class ShareLinkListView(LoginRequiredMixin, ListView):
    """分享链接列表"""
    model = ShareLink
    template_name = 'documents/share_link_list.html'
    context_object_name = 'share_links'
    paginate_by = 20
    
    def get_queryset(self):
        return ShareLink.objects.filter(
            created_by=self.request.user
        ).select_related('document').order_by('-created_at')


class DeleteShareLinkView(LoginRequiredMixin, DeleteView):
    """删除分享链接"""
    model = ShareLink
    template_name = 'documents/sharelink_confirm_delete.html'
    success_url = reverse_lazy('documents:share_link_list')
    
    def get_queryset(self):
        return ShareLink.objects.filter(created_by=self.request.user)


class DisableShareLinkView(LoginRequiredMixin, View):
    """禁用分享链接"""
    
    def post(self, request, pk):
        share_link = get_object_or_404(ShareLink, pk=pk, created_by=request.user)
        share_link.is_active = False
        share_link.save()
        messages.success(request, '分享链接已禁用')
        return redirect('documents:share_link_list')


class EnableShareLinkView(LoginRequiredMixin, View):
    """启用分享链接"""
    
    def post(self, request, pk):
        share_link = get_object_or_404(ShareLink, pk=pk, created_by=request.user)
        share_link.is_active = True
        share_link.save()
        messages.success(request, '分享链接已启用')
        return redirect('documents:share_link_list')


class DocumentVersionsView(LoginRequiredMixin, DetailView):
    """文档版本管理"""
    model = Document
    template_name = 'documents/document_versions.html'
    context_object_name = 'document'
    
    def get_queryset(self):
        return Document.objects.filter(author=self.request.user)
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        document = self.get_object()
        
        versions = DocumentVersion.objects.filter(
            document=document
        ).order_by('-created_at')
        
        context['versions'] = versions
        return context


class UploadVersionView(LoginRequiredMixin, View):
    """上传新版本"""
    
    def post(self, request, pk):
        document = get_object_or_404(Document, pk=pk, author=request.user)
        
        form = VersionForm(request.POST, request.FILES)
        if form.is_valid():
            with transaction.atomic():
                # 获取当前最新版本号
                latest_version = DocumentVersion.objects.filter(
                    document=document
                ).order_by('-created_at').first()
                
                if latest_version:
                    # 解析版本号并递增
                    # 支持 "v1" 或 "v1.0" 格式
                    version_str = latest_version.version_number[1:]  # 去掉 'v' 前缀
                    try:
                        # 尝试按小数点分割，提取主版本号
                        if '.' in version_str:
                            major_version = int(version_str.split('.')[0])
                        else:
                            major_version = int(version_str)
                        # 递增主版本号
                        new_major = major_version + 1
                        version_number = f"v{new_major}.0"
                    except (ValueError, IndexError):
                        # 如果解析失败，使用默认递增方式
                        try:
                            version_number = f"v{int(version_str) + 1}.0"
                        except ValueError:
                            version_number = "v2.0"
                else:
                    # 如果没有历史版本，说明当前文档是 v1.0，所以新版本应该是 v2.0
                    version_number = "v2.0"
                
                # 创建新版本
                version = form.save(commit=False)
                version.document = document
                version.version_number = version_number
                version.file_size = version.file.size
                version.created_by = request.user
                version.save()
                
                # 更新主文档文件：将版本文件复制到主文档位置，而不是直接指向版本文件
                # 这样主文档的文件路径保持不变，预览功能可以正常工作
                try:
                    old_file_path = document.file.path if document.file and os.path.exists(document.file.path) else None
                except:
                    old_file_path = None
                
                # 获取版本文件的路径
                try:
                    version_file_path = version.file.path
                except:
                    messages.error(request, '版本文件路径错误')
                    return redirect('documents:document_versions', pk=pk)
                
                # 如果版本文件存在，复制到主文档位置
                if os.path.exists(version_file_path):
                    from django.core.files import File
                    
                    # 打开版本文件并保存到主文档位置
                    with open(version_file_path, 'rb') as f:
                        # 保持主文档的文件名，只更新内容
                        file_name = os.path.basename(document.file.name) if document.file.name else os.path.basename(version.file.name)
                        document.file.save(file_name, File(f), save=False)
                    
                    document.file_size = version.file_size
                    document.save()
                    
                    # 删除旧文件（如果存在且不是版本文件）
                    if old_file_path and old_file_path != version_file_path and os.path.exists(old_file_path):
                        try:
                            # 检查旧文件是否被其他版本引用
                            other_versions = DocumentVersion.objects.filter(document=document).exclude(pk=version.pk)
                            is_referenced = False
                            for v in other_versions:
                                try:
                                    if v.file.path == old_file_path:
                                        is_referenced = True
                                        break
                                except:
                                    pass
                            
                            if not is_referenced:
                                os.remove(old_file_path)
                        except Exception:
                            pass  # 如果删除失败，忽略
                else:
                    messages.error(request, '版本文件不存在')
                    return redirect('documents:document_versions', pk=pk)
                
                messages.success(request, f'版本 {version_number} 上传成功')
        
        return redirect('documents:document_versions', pk=pk)


class DownloadVersionView(LoginRequiredMixin, View):
    """下载指定版本"""
    
    def get(self, request, pk, version_id):
        document = get_object_or_404(Document, pk=pk, author=request.user)
        version = get_object_or_404(DocumentVersion, pk=version_id, document=document)
        
        if not version.file or not os.path.exists(version.file.path):
            messages.error(request, '文件不存在')
            return redirect('documents:document_versions', pk=pk)
        
        response = FileResponse(
            version.file,
            content_type=mimetypes.guess_type(version.file.name)[0] or 'application/octet-stream'
        )
        # 强制下载，不预览
        response['Content-Disposition'] = f'attachment; filename="{document.title}_{version.version_number}{os.path.splitext(version.file.name)[1]}"'
        
        return response


class RestoreVersionView(LoginRequiredMixin, View):
    """恢复到指定版本"""
    
    def post(self, request, pk, version_id):
        document = get_object_or_404(Document, pk=pk, author=request.user)
        version = get_object_or_404(DocumentVersion, pk=version_id, document=document)
        
        with transaction.atomic():
            # 获取当前最新版本号，用于备份当前版本
            latest_version = DocumentVersion.objects.filter(
                document=document
            ).order_by('-created_at').first()
            
            # 生成备份版本号（基于最新版本号递增）
            if latest_version:
                version_str = latest_version.version_number[1:]  # 去掉 'v' 前缀
                try:
                    if '.' in version_str:
                        major_version = int(version_str.split('.')[0])
                    else:
                        major_version = int(version_str)
                    backup_version_number = f"v{major_version + 1}.0"
                except (ValueError, IndexError):
                    try:
                        backup_version_number = f"v{int(version_str) + 1}.0"
                    except ValueError:
                        backup_version_number = "v2.0"
            else:
                backup_version_number = "v2.0"
            
            # 备份当前版本
            current_version = DocumentVersion.objects.create(
                document=document,
                version_number=backup_version_number,
                file=document.file,
                file_size=document.file_size,
                change_log=f"恢复到 {version.version_number} 前的备份",
                created_by=request.user
            )
            
            # 恢复指定版本
            # 检查版本文件是否存在
            if version.file and os.path.exists(version.file.path):
                # 如果版本文件路径包含旧的 media/ 前缀，需要处理
                version_file_path = version.file.path
                # 如果路径中包含重复的 media，需要修正
                if 'media\\media\\' in version_file_path or 'media/media/' in version_file_path:
                    # 使用实际存在的文件路径
                    # 尝试去掉重复的 media
                    corrected_path = version_file_path.replace('media\\media\\', 'media\\').replace('media/media/', 'media/')
                    if os.path.exists(corrected_path):
                        version_file_path = corrected_path
                
                # 如果版本文件存在，将其复制到主文档位置（而不是直接指向版本文件）
                if os.path.exists(version_file_path):
                    from django.core.files import File
                    
                    # 保存旧文件路径（如果存在）
                    try:
                        old_file_path = document.file.path if document.file and os.path.exists(document.file.path) else None
                    except:
                        old_file_path = None
                    
                    # 将版本文件复制到主文档位置
                    with open(version_file_path, 'rb') as f:
                        # 保持主文档的文件名，只更新内容
                        file_name = os.path.basename(document.file.name) if document.file.name else os.path.basename(version.file.name)
                        document.file.save(file_name, File(f), save=False)
                    
                    document.file_size = version.file_size
                    document.save()
                    
                    # 注意：这里不删除旧文件，因为可能被其他版本引用
                else:
                    messages.error(request, f'版本文件不存在：{version_file_path}')
                    return redirect('documents:document_versions', pk=pk)
            else:
                messages.error(request, '版本文件不存在')
                return redirect('documents:document_versions', pk=pk)
            
            messages.success(request, f'已恢复到版本 {version.version_number}')
        
        return redirect('documents:document_versions', pk=pk)


# API视图
class DocumentInfoAPIView(LoginRequiredMixin, View):
    """文档信息API"""
    
    def get(self, request, pk):
        document = get_object_or_404(Document, pk=pk)
        
        # 权限检查
        if document.author != request.user and not document.is_public:
            return JsonResponse({'error': '无权限访问'}, status=403)
        
        return JsonResponse({
            'id': document.id,
            'title': document.title,
            'description': document.description,
            'author': document.author.get_full_name(),
            'file_size': document.file_size,
            'file_type': document.file_type,
            'download_count': document.download_count,
            'view_count': document.view_count,
            'created_at': document.created_at.isoformat(),
            'is_public': document.is_public,
            'can_edit': document.author == request.user,
        })


class CategoryListView(LoginRequiredMixin, ListView):
    """分类列表"""
    model = DocumentCategory
    template_name = 'documents/category_list.html'
    context_object_name = 'categories'
    paginate_by = 20
    
    def get_queryset(self):
        return DocumentCategory.objects.all().order_by('name')
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        # 获取所有分类
        all_categories = DocumentCategory.objects.all()
        
        # 区分管理员分类和用户分类
        if self.request.user.is_superuser:
            # 管理员可以看到所有分类
            admin_categories = all_categories.filter(created_by__is_superuser=True)
            user_categories = all_categories.filter(created_by__is_superuser=False)
        else:
            # 普通用户只能看到自己的分类
            admin_categories = all_categories.filter(created_by__is_superuser=True)
            user_categories = all_categories.filter(created_by=self.request.user)
        
        # 父分类和子分类
        context['admin_parent_categories'] = admin_categories.filter(parent__isnull=True)
        context['admin_child_categories'] = admin_categories.filter(parent__isnull=False)
        context['user_parent_categories'] = user_categories.filter(parent__isnull=True)
        context['user_child_categories'] = user_categories.filter(parent__isnull=False)
        
        # 统计信息
        context['admin_categories'] = admin_categories
        context['user_categories'] = user_categories
        context['active_categories'] = all_categories.filter(is_active=True)
        
        return context


class CreateCategoryView(LoginRequiredMixin, CreateView):
    """创建分类"""
    model = DocumentCategory
    form_class = CategoryForm
    template_name = 'documents/category_form.html'
    success_url = reverse_lazy('documents:category_list')
    
    def form_valid(self, form):
        form.instance.created_by = self.request.user
        return super().form_valid(form)


class EditCategoryView(LoginRequiredMixin, UpdateView):
    """编辑分类"""
    model = DocumentCategory
    form_class = CategoryForm
    template_name = 'documents/category_form.html'
    success_url = reverse_lazy('documents:category_list')
    
    def get_queryset(self):
        # 允许管理员编辑所有分类，普通用户只能编辑自己创建的分类
        if self.request.user.is_superuser:
            return DocumentCategory.objects.all()
        return DocumentCategory.objects.filter(created_by=self.request.user)


class DeleteCategoryView(LoginRequiredMixin, DeleteView):
    """删除分类"""
    model = DocumentCategory
    template_name = 'documents/category_confirm_delete.html'
    success_url = reverse_lazy('documents:category_list')
    
    def get_queryset(self):
        # 允许管理员删除所有分类，普通用户只能删除自己创建的分类
        if self.request.user.is_superuser:
            return DocumentCategory.objects.all()
        return DocumentCategory.objects.filter(created_by=self.request.user)


class UploadProgressAPIView(LoginRequiredMixin, View):
    """上传进度API"""
    def get(self, request):
        # 这里可以实现上传进度跟踪
        # 暂时返回简单的响应
        return JsonResponse({'progress': 0, 'status': 'pending'})


class CategoryDocumentsView(LoginRequiredMixin, ListView):
    """分类下的文档列表"""
    model = Document
    template_name = 'documents/category_documents.html'
    context_object_name = 'documents'
    paginate_by = 20
    
    def get_queryset(self):
        category_id = self.kwargs.get('pk')
        category = get_object_or_404(DocumentCategory, pk=category_id)
        
        # 获取该分类及其所有子分类
        subcategories = category.children.all()
        all_category_ids = [category.id] + list(subcategories.values_list('id', flat=True))
        
        # 查询该分类及其子分类下的所有文档
        queryset = Document.objects.filter(
            category__in=all_category_ids,
            author=self.request.user
        ).select_related('author', 'category').order_by('-created_at')
        
        return queryset
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        category_id = self.kwargs.get('pk')
        category = get_object_or_404(DocumentCategory, pk=category_id)
        
        context['category'] = category
        context['subcategories'] = category.children.all()
        
        return context


class DocumentPreviewView(LoginRequiredMixin, View):
    """文档预览视图"""
    
    def get(self, request, pk):
        document = get_object_or_404(Document, pk=pk)
        
        # 检查权限：文档作者、管理员或公开文档
        if not (document.author == request.user or 
                request.user.is_superuser or 
                request.user.is_admin() or 
                document.is_public):
            messages.error(request, '您没有权限预览此文档')
            return redirect('documents:document_list')
        
        # 记录查看次数
        document.view_count += 1
        document.save()
        
        # 获取文件路径，尝试多种可能的路径格式
        file_path = None
        try:
            original_path = document.file.path
            # 首先尝试原始路径
            if os.path.exists(original_path):
                file_path = original_path
            else:
                # 尝试修正路径（处理旧路径中重复的 media/）
                if 'media\\media\\' in original_path or 'media/media/' in original_path:
                    corrected_path = original_path.replace('media\\media\\', 'media\\').replace('media/media/', 'media/')
                    if os.path.exists(corrected_path):
                        file_path = corrected_path
                
                # 如果还是找不到，尝试使用 file.name 构建路径
                if not file_path and document.file.name:
                    file_name = document.file.name
                    # 去掉开头的 media/（如果存在）
                    if file_name.startswith('media/'):
                        file_name = file_name[6:]  # 去掉 'media/'
                    direct_path = os.path.join(settings.MEDIA_ROOT, file_name)
                    if os.path.exists(direct_path):
                        file_path = direct_path
                
                # 如果仍然找不到文件，记录错误但继续尝试预览（某些文件类型可能不需要本地路径）
                if not file_path:
                    # 对于某些文件类型（如PDF、图片），可以使用URL直接访问
                    # 对于需要本地路径的文件类型（如DOCX），会在后续处理中报错
                    pass
                    
        except (ValueError, OSError, FileNotFoundError) as e:
            # 如果获取路径时出错，记录但不立即返回
            # 某些文件类型可能不需要本地路径
            pass
        
        # 根据文件类型决定预览方式
        file_type = document.file_type.lower()
        
        if file_type in ['pdf']:
            # PDF文件直接在新窗口打开
            try:
                return redirect(document.file.url)
            except:
                messages.error(request, '无法访问文件')
                return redirect('documents:document_detail', pk=document.pk)
        elif file_type in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp']:
            # 图片文件显示预览页面
            return render(request, 'documents/document_preview.html', {
                'document': document,
                'preview_type': 'image',
                'file_url': document.file.url
            })
        elif file_type in ['docx']:
            # DOCX文件增强预览（需要本地路径）
            if not file_path or not os.path.exists(file_path):
                messages.error(request, '文件不存在，无法预览DOCX文件。请检查文件路径。')
                return redirect('documents:document_detail', pk=document.pk)
            try:
                from docx import Document as DocxDocument
                doc = DocxDocument(file_path)
                
                # 提取文档结构信息
                doc_info = {
                    'title': document.title,
                    'paragraphs': [],
                    'tables': [],
                    'images': [],
                    'total_paragraphs': 0,
                    'total_tables': 0,
                    'total_images': 0
                }
                
                # 提取段落和格式信息
                for i, paragraph in enumerate(doc.paragraphs):
                    if paragraph.text.strip():  # 只处理非空段落
                        # 安全地获取格式信息
                        para_bold = False
                        para_italic = False
                        para_underline = False
                        para_font_size = None
                        
                        try:
                            # 检查段落格式
                            for run in paragraph.runs:
                                if run.bold:
                                    para_bold = True
                                if run.italic:
                                    para_italic = True
                                if run.underline:
                                    para_underline = True
                                if run.font.size and not para_font_size:
                                    para_font_size = str(run.font.size.pt) + 'pt'
                        except AttributeError:
                            # 如果无法获取格式信息，使用默认值
                            pass
                        
                        para_info = {
                            'index': i + 1,
                            'text': paragraph.text,
                            'style': paragraph.style.name if hasattr(paragraph.style, 'name') else 'Normal',
                            'bold': para_bold,
                            'italic': para_italic,
                            'underline': para_underline,
                            'font_size': para_font_size,
                            'alignment': paragraph.alignment
                        }
                        
                        doc_info['paragraphs'].append(para_info)
                
                # 提取表格信息
                for table_idx, table in enumerate(doc.tables):
                    table_data = {
                        'index': table_idx + 1,
                        'rows': [],
                        'row_count': len(table.rows),
                        'col_count': len(table.columns) if table.rows else 0
                    }
                    
                    for row_idx, row in enumerate(table.rows):
                        row_data = {
                            'index': row_idx + 1,
                            'cells': []
                        }
                        for cell_idx, cell in enumerate(row.cells):
                            # 安全地获取单元格格式信息
                            cell_bold = False
                            cell_italic = False
                            
                            try:
                                # 尝试获取单元格的段落格式
                                for paragraph in cell.paragraphs:
                                    for run in paragraph.runs:
                                        if run.bold:
                                            cell_bold = True
                                        if run.italic:
                                            cell_italic = True
                            except AttributeError:
                                # 如果无法获取格式信息，使用默认值
                                pass
                            
                            cell_data = {
                                'index': cell_idx + 1,
                                'text': cell.text.strip(),
                                'bold': cell_bold,
                                'italic': cell_italic
                            }
                            row_data['cells'].append(cell_data)
                        table_data['rows'].append(row_data)
                    
                    doc_info['tables'].append(table_data)
                
                # 统计信息
                doc_info['total_paragraphs'] = len(doc_info['paragraphs'])
                doc_info['total_tables'] = len(doc_info['tables'])
                
                return render(request, 'documents/document_preview.html', {
                    'document': document,
                    'preview_type': 'docx_enhanced',
                    'doc_info': doc_info,
                    'file_url': document.file.url
                })
                
            except ImportError:
                messages.error(request, '系统缺少python-docx库，无法预览DOCX文件，请下载查看')
                return redirect('documents:document_detail', pk=document.pk)
            except Exception as e:
                messages.error(request, f'无法预览此DOCX文件：{str(e)}，请下载查看')
                return redirect('documents:document_detail', pk=document.pk)
        elif file_type in ['doc']:
            # DOC文件提示下载
            messages.info(request, 'DOC文件格式较旧，建议转换为DOCX格式后预览，或直接下载查看')
            return redirect('documents:document_detail', pk=document.pk)
        elif file_type in ['ppt', 'pptx']:
            # PPT/PPTX文件预览
            try:
                if file_type == 'pptx':
                    # 尝试使用python-pptx提取PPTX内容
                    try:
                        from pptx import Presentation
                        prs = Presentation(file_path)
                        
                        # 提取幻灯片信息
                        slides_info = []
                        for i, slide in enumerate(prs.slides):
                            slide_info = {
                                'slide_number': i + 1,
                                'shapes': [],
                                'notes': ''
                            }
                            
                            # 提取形状信息
                            for shape in slide.shapes:
                                if hasattr(shape, 'text') and shape.text.strip():
                                    slide_info['shapes'].append({
                                        'type': shape.shape_type,
                                        'text': shape.text.strip()[:200] + '...' if len(shape.text.strip()) > 200 else shape.text.strip()
                                    })
                            
                            # 提取备注
                            if slide.has_notes_slide:
                                notes_slide = slide.notes_slide
                                if notes_slide.notes_text_frame:
                                    slide_info['notes'] = notes_slide.notes_text_frame.text.strip()[:200] + '...' if len(notes_slide.notes_text_frame.text.strip()) > 200 else notes_slide.notes_text_frame.text.strip()
                            
                            slides_info.append(slide_info)
                        
                        return render(request, 'documents/document_preview.html', {
                            'document': document,
                            'preview_type': 'pptx',
                            'slides_info': slides_info,
                            'total_slides': len(slides_info),
                            'file_url': document.file.url
                        })
                    except ImportError:
                        # 如果没有安装python-pptx，显示基础信息
                        return render(request, 'documents/document_preview.html', {
                            'document': document,
                            'preview_type': 'ppt_basic',
                            'file_url': document.file.url
                        })
                else:
                    # PPT文件（旧格式）显示基础信息
                    return render(request, 'documents/document_preview.html', {
                        'document': document,
                        'preview_type': 'ppt_basic',
                        'file_url': document.file.url
                    })
            except Exception as e:
                messages.error(request, f'无法预览此PPT文件：{str(e)}，请下载查看')
                return redirect('documents:document_detail', pk=document.pk)
        elif file_type in ['txt', 'md', 'csv', 'json', 'xml', 'html', 'css', 'js', 'py', 'java', 'cpp', 'c']:
            # 文本文件显示内容（需要本地路径）
            if not file_path or not os.path.exists(file_path):
                messages.error(request, '文件不存在，无法预览文本文件。请检查文件路径。')
                return redirect('documents:document_detail', pk=document.pk)
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                return render(request, 'documents/document_preview.html', {
                    'document': document,
                    'preview_type': 'text',
                    'file_content': content,
                    'file_url': document.file.url
                })
            except UnicodeDecodeError:
                # 如果UTF-8解码失败，尝试其他编码
                try:
                    with open(file_path, 'r', encoding='gbk') as f:
                        content = f.read()
                    return render(request, 'documents/document_preview.html', {
                        'document': document,
                        'preview_type': 'text',
                        'file_content': content,
                        'file_url': document.file.url
                    })
                except:
                    messages.error(request, '无法预览此文件，请下载查看')
                    return redirect('documents:document_detail', pk=document.pk)
        else:
            # 不支持预览的文件类型
            messages.info(request, '此文件类型不支持在线预览，请下载查看')
            return redirect('documents:document_detail', pk=document.pk)


class DocumentReviewListView(LoginRequiredMixin, UserPassesTestMixin, ListView):
    """待审核文档列表（仅管理员可见）"""
    model = Document
    template_name = 'documents/document_review_list.html'
    context_object_name = 'documents'
    paginate_by = 20
    
    def test_func(self):
        return self.request.user.is_superuser or self.request.user.is_admin()
    
    def get_queryset(self):
        return Document.objects.filter(status='review').order_by('-created_at')


class DocumentReviewView(LoginRequiredMixin, UserPassesTestMixin, View):
    """文档审核视图（仅管理员可见）"""
    
    def test_func(self):
        return self.request.user.is_superuser or self.request.user.is_admin()
    
    def get(self, request, pk):
        document = get_object_or_404(Document, pk=pk, status='review')
        return render(request, 'documents/document_review.html', {
            'document': document
        })
    
    def post(self, request, pk):
        document = get_object_or_404(Document, pk=pk, status='review')
        action = request.POST.get('action')
        
        if action == 'approve':
            # 审核通过 - 状态变为已归档
            document.status = 'archived'
            document.save()
            messages.success(request, f'文档《{document.title}》审核通过，已归档')
            
            # 记录操作日志
            from system.models import SystemLog
            SystemLog.objects.create(
                level='INFO',
                message=f'管理员 {request.user.get_full_name()} 审核通过了文档《{document.title}》，已归档',
                module='documents',
                user=request.user,
                ip_address=request.META.get('REMOTE_ADDR')
            )
            
        elif action == 'reject':
            # 审核拒绝 - 状态变为审核未通过
            document.status = 'rejected'
            document.save()
            messages.success(request, f'文档《{document.title}》审核拒绝，教师可重新提交')
            
            # 记录操作日志
            from system.models import SystemLog
            SystemLog.objects.create(
                level='INFO',
                message=f'管理员 {request.user.get_full_name()} 拒绝了文档《{document.title}》，状态为审核未通过',
                module='documents',
                user=request.user,
                ip_address=request.META.get('REMOTE_ADDR')
            )
            
        
        return redirect('documents:document_review_list')